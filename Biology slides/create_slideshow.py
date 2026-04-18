from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Biology Slideshow: Immune and Integumentary Systems Interaction"
subtitle.text = "How they fight against germs and pathogens"

# Slide 2: Integumentary System
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Integumentary System'
tf = body_shape.text_frame
tf.text = 'The integumentary system includes the skin, hair, nails, and glands.'

p = tf.add_paragraph()
p.text = 'Functions: Protection, regulation of body temperature, sensation, excretion, vitamin D synthesis.'
p.level = 1

# Slide 3: Immune System
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Immune System'
tf = body_shape.text_frame
tf.text = 'The immune system defends the body against infections and diseases.'

p = tf.add_paragraph()
p.text = 'Components: Innate immunity (skin, mucous membranes) and adaptive immunity (T cells, B cells).'
p.level = 1

# Slide 4: Interaction between Systems
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Interaction between Immune and Integumentary Systems'
tf = body_shape.text_frame
tf.text = 'The skin acts as the first line of defense (physical barrier).'

p = tf.add_paragraph()
p.text = 'Immune cells in the skin (Langerhans cells) detect pathogens and initiate immune response.'
p.level = 1

p = tf.add_paragraph()
p.text = 'Wounds can allow pathogens to enter, triggering immune response.'
p.level = 1

# Slide 5: Fighting against Germs and Pathogens
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Fighting against Germs and Pathogens'
tf = body_shape.text_frame
tf.text = 'Physical barriers: Skin prevents entry of microbes.'

p = tf.add_paragraph()
p.text = 'Chemical barriers: Sweat, sebum kill bacteria.'
p.level = 1

p = tf.add_paragraph()
p.text = 'Immune response: Phagocytes engulf pathogens, antibodies neutralize toxins.'
p.level = 1

# Slide 6: Model of Fighting Pathogens
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = 'Model of How Systems Fight Pathogens'

# Create a simple diagram using matplotlib
fig, ax = plt.subplots(figsize=(8, 6))
ax.set_xlim(0, 10)
ax.set_ylim(0, 10)
ax.axis('off')

# Draw boxes
rect1 = patches.FancyBboxPatch((1, 7), 2, 1, boxstyle="round,pad=0.1", facecolor="lightblue", edgecolor="black")
ax.add_patch(rect1)
ax.text(2, 7.5, 'Skin Barrier', ha='center', va='center', fontsize=12)

rect2 = patches.FancyBboxPatch((4, 7), 2, 1, boxstyle="round,pad=0.1", facecolor="lightgreen", edgecolor="black")
ax.add_patch(rect2)
ax.text(5, 7.5, 'Pathogen Entry', ha='center', va='center', fontsize=12)

rect3 = patches.FancyBboxPatch((7, 7), 2, 1, boxstyle="round,pad=0.1", facecolor="lightcoral", edgecolor="black")
ax.add_patch(rect3)
ax.text(8, 7.5, 'Immune Response', ha='center', va='center', fontsize=12)

# Arrows
ax.arrow(3, 7.5, 1, 0, head_width=0.1, head_length=0.1, fc='black', ec='black')
ax.arrow(6, 7.5, 1, 0, head_width=0.1, head_length=0.1, fc='black', ec='black')

# Save the figure
plt.savefig('model_diagram.png', bbox_inches='tight', dpi=150)
plt.close()

# Add the image to the slide
left = Inches(1)
top = Inches(1.5)
height = Inches(5)
pic = slide.shapes.add_picture('model_diagram.png', left, top, height=height)

# Slide 7: Conclusion
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Conclusion'
tf = body_shape.text_frame
tf.text = 'The integumentary and immune systems work together to protect the body from pathogens.'

p = tf.add_paragraph()
p.text = 'The skin provides a physical barrier, while the immune system provides cellular and humoral defenses.'
p.level = 1

# Save the presentation
prs.save('biology_slideshow.pptx')
print("Slideshow created successfully: biology_slideshow.pptx")